import os
import argparse
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO

def merge_individual_style(merged_pptx, current_pptx):
    """모드 1: 원본 내용과 디자인을 최대한 유지하며 합치기"""
    for slide in current_pptx.slides:
        # 원본 슬라이드의 레이아웃을 그대로 복제 시도
        try:
            slide_layout = slide.slide_layout
            new_slide = merged_pptx.slides.add_slide(slide_layout)
            for shape in slide.shapes:
                el = shape.element
                new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')
        except Exception as e:
            print(f"슬라이드 복제 중 오류 발생: {e}")

def merge_unified_theme(merged_pptx, current_pptx):
    """모드 2: 특정 파일(첫 번째 파일)의 스타일을 전체에 강제 적용"""
    # 기준 파일에서 사용 가능한 레이아웃 확인 (IndexError 방지)
    layouts = merged_pptx.slide_layouts
    # 보통 1번(제목 및 내용)을 쓰되, 없으면 0번(제목)을 사용
    target_layout = layouts[1] if len(layouts) > 1 else layouts[0]

    for slide in current_pptx.slides:
        new_slide = merged_pptx.slides.add_slide(target_layout)
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_stream = BytesIO(shape.image.blob)
                new_slide.shapes.add_picture(img_stream, shape.left, shape.top, shape.width, shape.height)
            elif shape.has_text_frame:
                # 텍스트 박스 생성 (위치 유지)
                new_shape = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                new_shape.text_frame.text = shape.text_frame.text

def main():
    parser = argparse.ArgumentParser(description="PPT 병합 전문 도구")
    parser.add_argument('-d', '--dir', type=str, required=True, help='PPT 폴더 경로')
    parser.add_argument('-o', '--out', type=str, default='merged_result.pptx', help='결과 파일명')
    parser.add_argument('-m', '--mode', type=str, default='theme', choices=['theme', 'individual'],
                        help='theme: 스타일 통일, individual: 원본 유지')

    args = parser.parse_args()

    files = sorted([f for f in os.listdir(args.dir) if f.endswith('.pptx')])
    if not files: return print("파일이 없습니다.")

    # 첫 번째 파일을 마스터/베이스로 지정
    merged_pptx = Presentation(os.path.join(args.dir, files[0]))
    
    for file in files[1:]:
        print(f"처리 중: {file} (모드: {args.mode})")
        current_pptx = Presentation(os.path.join(args.dir, file))
        
        if args.mode == 'individual':
            merge_individual_style(merged_pptx, current_pptx)
        else:
            merge_unified_theme(merged_pptx, current_pptx)

    merged_pptx.save(args.out)
    print(f"병합 완료: {args.out}")

if __name__ == "__main__":
    main()